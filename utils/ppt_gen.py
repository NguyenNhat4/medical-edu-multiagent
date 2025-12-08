from pptx import Presentation
from pptx.util import Inches, Pt

def generate_ppt(slides_data, filename):
    """
    Generates a PowerPoint presentation from a list of slide data.

    Args:
        slides_data (list): List of dicts, each containing 'title', 'content' (list), and 'speaker_notes'.
        filename (str): Output filename.
    """
    prs = Presentation()

    # Title Slide (First slide usually)
    # We can assume the first item is the title slide if we want, or just generic formatting.
    # For this simple agent, we'll use Layout 1 (Title and Content) for all,
    # unless we detect it's the very first one which might be Title Only.

    # Let's just use Title and Content for all for simplicity.

    for i, slide_info in enumerate(slides_data):
        # Use layout 1: Title and Content
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # Set Title
        title = slide.shapes.title
        title.text = slide_info.get("title", "Untitled")

        # Set Content (Bullet points)
        # Check if placeholders[1] exists (Body)
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.word_wrap = True

            content_points = slide_info.get("content", [])
            if isinstance(content_points, str):
                content_points = [content_points]

            for idx, point in enumerate(content_points):
                if idx == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = point
                p.font.size = Pt(20)

        # Add Speaker Notes
        notes_text = slide_info.get("speaker_notes", "")
        if notes_text:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes_text

    prs.save(filename)
    return filename

if __name__ == "__main__":
    # Test
    sample_data = [
        {
            "title": "Slide 1",
            "content": ["Point A", "Point B"],
            "speaker_notes": "Note for slide 1"
        },
        {
            "title": "Slide 2",
            "content": ["Point C", "Point D"],
            "speaker_notes": "Note for slide 2"
        }
    ]
    generate_ppt(sample_data, "test.pptx")
    print("Generated test.pptx")
